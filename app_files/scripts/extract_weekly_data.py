import os, re, subprocess, json, zipfile, xml.etree.ElementTree as ET
from pathlib import Path
from datetime import datetime

DATA_DIR = Path(os.environ.get('BOP_REPORTS_DIR', str(Path(__file__).resolve().parents[1] / 'reports')))
files = sorted([p for p in DATA_DIR.iterdir() if p.is_file() and p.suffix.lower() in ['.pdf','.pptx']])

# Basic helpers

def pptx_text(path: Path) -> str:
    # use python-pptx if available; fallback to XML extraction
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slide_texts=[]
        for i, slide in enumerate(prs.slides, start=1):
            parts=[f"<SLIDE {i}>"]
            for shape in slide.shapes:
                try:
                    if getattr(shape, 'has_table', False):
                        for row in shape.table.rows:
                            parts.append(' | '.join((cell.text or '').replace('\n',' ').strip() for cell in row.cells))
                except Exception:
                    pass
                if hasattr(shape, 'text') and shape.text:
                    parts.append(shape.text)
            slide_texts.append('\n'.join(parts))
        return '\n\n'.join(slide_texts)
    except Exception:
        texts=[]
        ns={'a':'http://schemas.openxmlformats.org/drawingml/2006/main'}
        with zipfile.ZipFile(path) as z:
            slide_names=sorted([n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')], key=lambda x:int(re.search(r'slide(\d+)\.xml',x).group(1)))
            for idx,name in enumerate(slide_names, start=1):
                xml=z.read(name)
                root=ET.fromstring(xml)
                texts.append(f"<SLIDE {idx}>\n"+'\n'.join(t.text or '' for t in root.findall('.//a:t', ns)))
        return '\n\n'.join(texts)

def pdf_text(path: Path) -> str:
    # Prefer poppler pdftotext when available because it preserves table layout well.
    # Fall back to pypdf for Windows users who do not have poppler installed.
    try:
        res = subprocess.run(['pdftotext','-layout',str(path),'-'], capture_output=True, text=True, timeout=60)
        if res.returncode == 0 and res.stdout.strip():
            return res.stdout
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for idx, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ''
            pages.append(f'<PAGE {idx}>\n{text}')
        return '\n\n'.join(pages)
    except Exception:
        return ''

def extract_text(path: Path) -> str:
    if path.suffix.lower()=='.pdf':
        return pdf_text(path)
    elif path.suffix.lower()=='.pptx':
        txt = pptx_text(path)
        # Some weekly PPTX files contain schedule tables as rendered objects.
        # If BOP Distribution  is present but table values are not exposed,
        # try a LibreOffice-to-PDF fallback when soffice/libreoffice exists.
        if 'BOP | Distribution ()' in txt and ' (Schedule)' not in txt:
            try:
                import tempfile, shutil
                soffice = shutil.which('libreoffice') or shutil.which('soffice')
                if soffice:
                    with tempfile.TemporaryDirectory(prefix='bop_pptx_pdf_') as tmp:
                        subprocess.run([soffice, '--headless', '--convert-to', 'pdf', '--outdir', tmp, str(path)], capture_output=True, text=True, timeout=120)
                        pdf = Path(tmp) / (path.stem + '.pdf')
                        if pdf.exists():
                            txt_pdf = pdf_text(pdf)
                            if txt_pdf.strip():
                                txt = txt + '\n\n' + txt_pdf
            except Exception:
                pass
        return txt
    return ''

def to_float(s):
    if s is None: return None
    s=s.strip().replace('%','').replace(' ','')
    if not s: return None
    # convert comma decimal, preserve thousands dots loosely
    if ',' in s and '.' in s:
        s=s.replace('.','').replace(',','.')
    else:
        s=s.replace(',','.')
    try:
        return float(s)
    except Exception:
        return None

def to_int_value(s):
    if s is None: return None
    # if value like 16/17, return first and last separately elsewhere
    m=re.search(r'\d+',str(s))
    return int(m.group(0)) if m else None

def parse_pair(s):
    if not s: return (None,None)
    parts=str(s).split('/')
    if len(parts)>=2:
        return to_int_value(parts[0]), to_int_value(parts[1])
    return to_int_value(s), None

def normalize(s):
    return re.sub(r'\s+', ' ', s or '').strip()

def parse_hse(text, fname):
    rec = {
        'file_name': fname,
        'week': None,
        'issue_date': None,
        'lti_free_days': None,
        'lti_free_label': None,
        'ltifr': None,
        'trifr': None,
        'site_walks': None, 'site_walks_last': None, 'site_walks_ptd': None,
        'observations': None, 'observations_last': None, 'observations_ptd': None,
        'bbs': None, 'bbs_last': None, 'bbs_ptd': None,
        'target_inspections': None, 'target_inspections_last': None, 'target_inspections_ptd': None,
        'near_miss': None, 'near_miss_last': None, 'near_miss_ptd': None,
        'permits_to_work': None, 'permits_last': None, 'permits_ptd': None,
        'rams_approved': None, 'rams_last': None, 'rams_ptd': None,
        'rewards': None, 'rewards_last': None, 'rewards_ptd': None,
        'sending_offs': None, 'sending_last': None, 'sending_ptd': None,
        'positive_pct': None, 'negative_pct': None,
        'incident_summary': None,
        'highlights': None,
        'top_obs_category': None,
    }
    m=re.search(r'Weekly Report\s+W\s*(\d{1,2})', text, re.I)
    if not m:
        m=re.search(r'(?:Summary|Dashboard)\s*\(Week\s*(\d{1,2})\)', text, re.I)
    if m: rec['week']=int(m.group(1))
    m=re.search(r'Issue date[: ]+([^\n\r]+)', text, re.I)
    if m: rec['issue_date']=normalize(m.group(1))
    m=re.search(r'(\d+)\s+DAYS\s*-\s*LTI FREE', text, re.I)
    if m:
        rec['lti_free_days']=int(m.group(1)); rec['lti_free_label']=m.group(0).replace('\n',' ')
    elif re.search(r'1\s+YEAR\s*-\s*LTI\s*FREE', text, re.I):
        rec['lti_free_days']=365; rec['lti_free_label']='1 YEAR - LTI FREE'
    # HSE table patterns
    metrics = [
        ('SITE WALKS','site_walks'), ('OBSERVATIONS','observations'), ('BBS INSPECTIONS','bbs'), ('BBS ASSESSMENTS','bbs'),
        ('TARGET INSPECTIONS','target_inspections'), ('NEAR-MISS','near_miss'), ('NEAR MISS','near_miss'),
        ('PERMITS TO WORK','permits_to_work'), ('RAMS APPROVED','rams_approved'), ('REWARDS','rewards'), ('SENDING OFFS','sending_offs')]
    # parse table rows with week, last week, project to date
    for label, key in metrics:
        # Allow whitespace and hyphen variation
        pat_label = label.replace(' ', r'\s+').replace('-', r'[-\s]?')
        # need ensure not matching in prose; expect three integers after label
        mm = re.search(pat_label + r'\s*(?:\|\s*)?([0-9]+)\s*(?:\|\s*)?([0-9]+)\s*(?:\|\s*)?([0-9]+)', text, re.I)
        if mm:
            rec[key]=int(mm.group(1)); rec[key+'_last' if key not in ['sending_offs'] else 'sending_last']=int(mm.group(2)); rec[key+'_ptd' if key not in ['sending_offs'] else 'sending_ptd']=int(mm.group(3))
    # parse older leading block (W05-W14) and line-level text where safety-moment text appears between KPI lines
    if rec['site_walks'] is None:
        lead_window = text[:2500]
        old_fields = [
            ('site_walks', r'(^|\n)[^\n]*?([0-9/]+)[ \t]+Site Walks'),
            ('observations', r'(^|\n)[^\n]*?([0-9/]+)[ \t]+Observations'),
            ('bbs', r'(^|\n)[^\n]*?([0-9/]+)[ \t]+BBS[ \t]+(?:Assessments|Inspections)'),
            ('target_inspections', r'(^|\n)[^\n]*?([0-9/]+)[ \t]+Target Inspections'),
            ('near_miss', r'(^|\n)[^\n]*?([0-9/]+)[ \t]+Near Miss'),
        ]
        for key, pat in old_fields:
            mm=re.search(pat, lead_window, re.I)
            if mm:
                val = mm.group(2) if mm.lastindex and mm.lastindex >= 2 else mm.group(1)
                v,last=parse_pair(val)
                rec[key]=v
                if last is not None:
                    rec[key+'_last']=last
    # LTIFR/TRIFR - table or formulas
    # take first sensible decimal with comma/dot after LTIFR/TRIFR
    m=re.search(r'LTIFR\s*(?:\|\s*)?([0-9]+[,.][0-9]+)', text, re.I)
    if m: rec['ltifr']=to_float(m.group(1))
    m=re.search(r'TRIFR?\s*(?:\|\s*)?([0-9]+[,.][0-9]+)', text, re.I)
    if m: rec['trifr']=to_float(m.group(1))
    # formula variants
    m_all=re.findall(r'LTIFR[^\n=]*=.*?=\s*([0-9]+[,.][0-9]+)', text, re.I)
    if m_all: rec['ltifr']=to_float(m_all[-1])
    m_all=re.findall(r'TRIF[^\n=]*=.*?=\s*([0-9]+[,.][0-9]+)', text, re.I)
    if m_all: rec['trifr']=to_float(m_all[-1])
    # Observation positive/negative in dashboard
    m=re.search(r'Positive\s+(\d+)\s*%', text, re.I)
    if m: rec['positive_pct']=int(m.group(1))
    m=re.search(r'Negative\s+(\d+)\s*%', text, re.I)
    if m: rec['negative_pct']=int(m.group(1))
    # Incident/announcement summary
    for marker in ['Announcements:', 'Incident Announcements', 'Highlights/Downlights of the week:', 'Highlights / Lowlights', 'Highlights:']:
        idx=text.lower().find(marker.lower())
        if idx!=-1:
            snippet=text[idx:idx+1200]
            # stop at issue date or page break perhaps after first block
            snippet=re.split(r'\n\s*(?:BOP\||EHSS|BOP-IDOM|Issue date|CONFIDENTIAL|Strictly confidential)\b', snippet, maxsplit=1)[0]
            snippet=normalize(snippet.replace(marker,''))
            if snippet:
                if 'Highlight' in marker:
                    rec['highlights']=snippet[:500]
                else:
                    rec['incident_summary']=snippet[:500]
                break
    # Observation category: find highest explicit category near Observations by Category
    m=re.search(r'Observations by Category(.{0,400})', text, re.I|re.S)
    if m:
        cats=[]
        for line in m.group(1).splitlines():
            line=normalize(line)
            mm=re.match(r'([A-Za-z][A-Za-z /&\.\-]+?)\s+(\d{1,3})\b', line)
            if mm:
                cats.append((mm.group(1).strip(), int(mm.group(2))))
        if cats:
            cats=sorted(cats, key=lambda x:x[1], reverse=True)
            rec['top_obs_category']=f"{cats[0][0]} ({cats[0][1]})"
    # normalize alternate dynamic keys from generic parsing
    if rec.get('permits_to_work_last') is not None:
        rec['permits_last']=rec.get('permits_to_work_last')
    if rec.get('permits_to_work_ptd') is not None:
        rec['permits_ptd']=rec.get('permits_to_work_ptd')
    if rec.get('rams_approved_last') is not None:
        rec['rams_last']=rec.get('rams_approved_last')
    if rec.get('rams_approved_ptd') is not None:
        rec['rams_ptd']=rec.get('rams_approved_ptd')
    return rec

# schedule and area data from latest / any report

def parse_overall_schedule(text, fname):
    rec={'file_name':fname,'week':None,'cutoff':None,'actual_pct':None,'baseline_pct':None,'forecast_pct':None,'deviation_pct':None,
         'engineering_cumm_actual':None,'engineering_cumm_plan':None,'engineering_cumm_forecast':None,'engineering_deviation_pct':None,
         'procurement_cumm_actual':None,'procurement_cumm_plan':None,'procurement_cumm_forecast':None,'procurement_deviation_pct':None,
         'construction_cumm_actual':None,'construction_cumm_plan':None,'construction_cumm_forecast':None,'construction_deviation_pct':None,
         'overall_week_actual':None}
    m=re.search(r'Weekly Report\s+W\s*(\d{1,2})', text, re.I)
    if m: rec['week']=int(m.group(1))
    m=re.search(r'Overall Cut-?off:\s*([^\n\r]+)', text, re.I)
    if m: rec['cutoff']=normalize(m.group(1))
    for key,pat in [('actual_pct',r'Actual:\s*([0-9]+[,.][0-9]+)%'),('baseline_pct',r'BL\s*CCE\d+:\s*([0-9]+[,.][0-9]+)%'),('deviation_pct',r'Deviation:\s*([-]?[0-9]+[,.][0-9]+)%'),('forecast_pct',r'Forecast:\s*([0-9]+[,.][0-9]+)%')]:
        m=re.search(pat, text, re.I)
        if m: rec[key]=to_float(m.group(1))

    # Parse discipline-level cumulative rows from the Overall Progress table.
    # The reports have wide tables; pdftotext -layout preserves the rows but the
    # "Overall Cut-off" marker may be more than 3000 chars away, so use a wider
    # bounded window and read the first ACT/Plan/FC percentages after Cumm.
    block_match=re.search(r'Overall Progress(.{0,7000}?)Overall Cut-?off:', text, re.I|re.S)
    if not block_match:
        block_match=re.search(r'Overall Progress(.{0,7000}?)(?:Main deviations|6-month key Waypoints|BOP \||Status:)', text, re.I|re.S)
    if block_match:
        block=block_match.group(1)
        mapping=[('Engineering','engineering'),('Procurement','procurement'),('Construction','construction')]
        lines=block.splitlines()
        for discipline, prefix in mapping:
            values=None
            for i,line in enumerate(lines):
                # The discipline label is usually isolated on the left side of the table.
                if not re.search(r'\b'+re.escape(discipline)+r'\b', line, re.I):
                    continue
                for look in lines[i+1:i+7]:
                    pcts=re.findall(r'([0-9]+[,.][0-9]+)%', look)
                    nums=[to_float(x) for x in pcts]
                    nums=[x for x in nums if x is not None]
                    # Cumulative rows have large percent values; weekly rows are usually below 5%.
                    if len(nums) >= 3 and nums[0] >= 10:
                        values=nums[:3]
                        break
                if values:
                    break
            if not values:
                dm=re.search(discipline+r'.{0,1200}?Cumm\.\s*([0-9]+[,.][0-9]+)%\s+([0-9]+[,.][0-9]+)%\s+([0-9]+[,.][0-9]+)%', block, re.I|re.S)
                if dm:
                    values=[to_float(dm.group(1)), to_float(dm.group(2)), to_float(dm.group(3))]
            if values:
                actual, plan, fc = values[0], values[1], values[2]
                rec[f'{prefix}_cumm_actual']=actual
                rec[f'{prefix}_cumm_plan']=plan
                rec[f'{prefix}_cumm_forecast']=fc
                if actual is not None and fc is not None:
                    rec[f'{prefix}_deviation_pct']=round(actual-fc, 2)
        # The Overall row follows the Construction row in current templates.
        om=re.search(r'Overall\s+.*?Week\s+([0-9]+[,.][0-9]+)%', block, re.I|re.S)
        if om: rec['overall_week_actual']=to_float(om.group(1))

    # Fallback for pages titled BOP | Engineering, which are sometimes easier to
    # parse than the wide overall table. Keep only if the overall table did not provide a value.
    if rec.get('engineering_cumm_actual') is None:
        em=re.search(r'BOP \| Engineering.{0,500}?Actual progress:\s*([0-9]+[,.][0-9]+)%.{0,120}?Forecast Progress:\s*([0-9]+[,.][0-9]+)%', text, re.I|re.S)
        if em:
            actual=to_float(em.group(1)); fc=to_float(em.group(2))
            rec['engineering_cumm_actual']=actual
            rec['engineering_cumm_forecast']=fc
            if actual is not None and fc is not None:
                rec['engineering_deviation_pct']=round(actual-fc, 2)
    return rec


def parse_waypoints(text, fname):
    rows=[]
    week=None
    m=re.search(r'Weekly Report\s+W\s*(\d{1,2})', text, re.I)
    if m: week=int(m.group(1))
    # Target main block 6-month key waypoints around Overall Progress
    # We'll scrape all lines containing dates in Key Waypoints blocks but limit to known labels.
    known=[
        'Temporary Water (Ready to Supply)',
        'Temporary Power (Ready to Operate)',
        'Temporary Power',
        'Power Supply from NPB-511 to ASP',
        'Piping Connections - Priority 1 Finish',
        'P1 – Pipebridge Installed + Piping Connections Finish',
        'P2 – Pipebridge Installed + Piping Connections Finish',
        'P5 – Pipebridge Installed + Piping Connections Finish',
        'P4 – Pipebridge Installed + Piping Connections Finish',
        'P2 – Pipebridge',
        'P5 – Pipebridge',
        'P4 – Pipebridge',
        'Priority 1 UG piping installation completion',
        'Route 1 UG pipe completion',
        'RW/FF Pumps area UG piping works Start',
        'RW/FF Pumps area',
        'Electrical Works by STS Start',
        'Start Electrical Works',
        'East water pond 1,2 & 3 start',
        'East water pond 1, 2 & 3 start',
        'Railway Groundworks Crossings Start',
        'Raw Water Supply',
    ]
    # Collapse blocks preserving lines
    for label in known:
        pat=re.escape(label).replace('\\–','[–-]').replace('\\-','[-–]')
        mm=re.search(pat+r'.{0,120}?([0-9]{1,2}\s+[A-Za-z]{3,9}\s+\d{4})\s+([0-9]{1,2}\s+[A-Za-z]{3,9}\s+\d{4}|TBC)', text, re.I|re.S)
        if mm:
            planned, forecast = mm.group(1), mm.group(2)
            duplicate = any(
                r['planned_date'] == planned
                and r['forecast_date'] == forecast
                and (label in r['waypoint'] or r['waypoint'] in label)
                for r in rows
            )
            if not duplicate:
                rows.append({'week':week,'file_name':fname,'area':'Overall','waypoint':label,'planned_date':planned,'forecast_date':forecast})
    return rows

# Extract latest area statuses and watchlist from W24 (or latest available)

def add_watch(rows, week, source_file, area, issue, category='Schedule/Interface', severity='High', owner='', action='', status='Open', due=''):
    rows.append({'week':week,'file_name':source_file,'area':area,'category':category,'severity':severity,'issue':normalize(issue),'owner':owner,'action':normalize(action),'status':status,'due_date':due})

# Parse latest known areas from text, with special regex blocks for W24.
def parse_area_statuses(text, fname):
    week=None
    m=re.search(r'Weekly Report\s+W\s*(\d{1,2})', text, re.I)
    if m: week=int(m.group(1))
    rows=[]
    # Helper for management sections
    section_defs=[
        ('Rail', r'BOP \| Railways Management.*?(?=\n\s*BOP \| Roads|\Z)'),
        ('Roads', r'BOP \| Roads & Ditches Management.*?(?=\n\s*BOP \| Ponds|\Z)'),
        ('Ponds', r'BOP \| Ponds Management.*?(?=\n\s*BOP \| (?!Ponds)|\Z)'),
        ('NPB/Buildings', r'BOP \| Buildings & Other IDOM.*?(?=\n\s*BOP \| Engineering|\Z)'),
    ]
    for area, pat in section_defs:
        mm=re.search(pat, text, re.I|re.S)
        if not mm: continue
        block=mm.group(0)
        actual=forecast=deviation=None
        am=re.search(r'Actual progress:\s*([0-9]+[,.][0-9]+)%', block, re.I)
        if am: actual=to_float(am.group(1))
        fm=re.search(r'Forecast progress:\s*([0-9]+[,.][0-9]+)%', block, re.I)
        if fm: forecast=to_float(fm.group(1))
        dm=re.search(r'Deviation[: ]\s*([-]?[0-9]+(?:[,.][0-9]+)?)%', block, re.I)
        if dm: deviation=to_float(dm.group(1))
        # Area status row as issue
        if actual is not None:
            add_watch(rows, week, fname, area, f"Latest progress: actual {actual:.2f}%, forecast {forecast if forecast is not None else 'n/a'}%, deviation {deviation if deviation is not None else 'n/a'}%", category='Progress', severity='Info', status='Monitor')
        # 1+6 activities
        acts=re.findall(r'Week\s+(\d{2})\s*[-–]\s*([^\n]+)', block)
        for wk, act in acts[:8]:
            add_watch(rows, week, fname, area, f"W{wk} look-ahead: {act}", category='Look-ahead', severity='Medium', status='Planned')
        # red activities
        for red, mitigation in re.findall(r'\n\s*(?:PRC|CON|Eng|Con|Proc|Discipline)?\s{1,}([^\n]{10,100}?)\s{2,}([^\n]{10,200})', block):
            # screen out headers and waypoints
            if any(x in red.lower() for x in ['key waypoints','planned date','cut-off','actual progress','forecast progress','week ']):
                continue
            if any(x in mitigation.lower() for x in ['planned date','forecast date','achievement']):
                continue
            if len(rows) < 200:
                add_watch(rows, week, fname, area, red, action=mitigation, category='Red activity / mitigation', severity='High')
    # Keyword-based issues from latest quality/risk text for area responsibilities
    keywords={
        'Rail':['Rail','railway','Rails'],
        'Ponds':['Pond','PONDS','West Pond','Contact water pond','East stormwater pond'],
        'Roads':['Road','Roads','Avenue 9.5','Av 9.5','Road 6'],
        'NPB/Buildings':['NPB 100','NPB 200','NPB 103','NPB 202','NPB 203','NPB 205','NPB 206','NPB 207','NPB102','NPB205']
    }
    # split into sentences/lines and capture relevant lines
    lines=[normalize(l) for l in text.splitlines()]
    for area,kws in keywords.items():
        for line in lines:
            if len(line)<20 or len(line)>300: continue
            if any(k.lower() in line.lower() for k in kws):
                if any(no in line.lower() for no in ['weekly report','issue date','confidential','last update']): continue
                # classify severity based on words
                sev='Medium'
                if re.search(r'pending|delay|risk|without|rejected|missing|not approved|cannot|blocked|hold|overdue|red|damage|leak', line, re.I): sev='High'
                cat='Keyword watch'
                add_watch(rows, week, fname, area, line, category=cat, severity=sev, status='Open' if sev=='High' else 'Monitor')
    # Deduplicate
    seen=set(); out=[]
    for r in rows:
        key=(r['area'], r['issue'][:120], r['week'])
        if key in seen: continue
        seen.add(key); out.append(r)
    return out


def _month_to_date(text, prefer_end=True):
    if not text:
        return None
    import calendar
    raw=normalize(str(text).replace('TB C','TBC'))
    if raw.upper() == 'TBC':
        return None
    raw=re.sub(r'\bTBC\b','', raw, flags=re.I).strip()
    months={m.lower():i for i,m in enumerate(['Jan','Feb','Mar','Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec'],1)}
    m=re.search(r'(\d{1,2})\s+([A-Za-z]{3,9})\s+(\d{2,4})', raw)
    if m:
        d=int(m.group(1)); mon=m.group(2)[:3].lower(); y=int(m.group(3)); y=2000+y if y<100 else y
        return f"{y:04d}-{months.get(mon,1):02d}-{d:02d}"
    m=re.search(r'([A-Za-z]{3,9})\s+(\d{2,4})', raw)
    if m:
        mon=m.group(1)[:3].lower(); y=int(m.group(2)); y=2000+y if y<100 else y
        mm=months.get(mon,1); d=calendar.monthrange(y,mm)[1] if prefer_end else 1
        return f"{y:04d}-{mm:02d}-{d:02d}"
    return None


def _pct_values_after(block, label, stop=None):
    idx=block.lower().find(label.lower())
    if idx < 0:
        return {}
    end=len(block)
    if stop:
        j=block.lower().find(stop.lower(), idx+1)
        if j > 0:
            end=j
    seg=block[idx:end]
    pcts=re.findall(r'([0-9]+[,.][0-9]+)%|([0-9]{1,3})%', seg)
    nums=[]
    for a,b in pcts:
        v=to_float(a or b)
        if v is not None:
            nums.append(v)
    if len(nums) < 6:
        return {}
    nums=nums[:6]
    return {'actual':nums[0], 'plan':nums[1], 'forecast':nums[2], 'ach_act_fc':nums[3], 'ach_act_plan':nums[4], 'next_fc':nums[5]}


def parse_ugp_distribution(text, fname):
    week=None
    m=re.search(r'Weekly Report\s+W\s*(\d{1,2})', text, re.I)
    if m:
        week=int(m.group(1))
    issue_date=None
    m=re.search(r'Issue date[: ]+([^\n\r]+)', text, re.I)
    if m:
        issue_date=normalize(m.group(1))
    m=re.search(r'BOP \| Distribution \(\).*?(?=\n\s*BOP \| Buildings|\Z)', text, re.I|re.S)
    if not m:
        return None
    block=normalize(m.group(0).replace('TB C','TBC'))
    sched=_pct_values_after(block, ' (Schedule)', 'Main deviations')
    if not sched:
        return None
    eng=_pct_values_after(block, 'Engineering (MDR)', 'Procurement')
    proc=_pct_values_after(block, 'Procurement (TG)', '')
    route_names=[
        ' – Route 3 Commencement', ' – Route 0 testing complete (After winter)', ' – Route 0 testing complete',
        ' – Route 5 Completion', ' – Route 2 Commencement', ' – RW & FW Completion',
        ' – Route 1 Completion (incl. Rd. 9.5)', ' – Route 4 Commencement (Rd. 9.4)', ' – Route 4 Commencement'
    ]
    idxs=[]
    for name in route_names:
        i=block.find(name)
        if i >= 0:
            idxs.append((i,name))
    idxs=sorted(idxs)
    unique=[]
    for i,n in idxs:
        if any(abs(i-j)<5 for j,_ in unique):
            continue
        unique.append((i,n))
    waypoints=[]
    for k,(i,name) in enumerate(unique):
        end=unique[k+1][0] if k+1 < len(unique) else block.lower().find('1. achievement', i)
        if end < 0:
            end=min(len(block), i+180)
        seg=block[i+len(name):end]
        dts=re.findall(r'(?:\d{1,2}\s+[A-Za-z]{3,9}\s+\d{2,4}|[A-Za-z]{3,9}\s+\d{2,4}\s*TBC|[A-Za-z]{3,9}\s+\d{2,4}|TBC)', seg)
        planned=normalize(dts[0]) if len(dts) >= 1 else ''
        forecast=normalize(dts[1]) if len(dts) >= 2 else (normalize(dts[0]) if len(dts)==1 else '')
        waypoints.append({
            'name': name.replace(' – ','').strip(), 'planned_raw': planned, 'forecast_raw': forecast,
            'planned': _month_to_date(planned, prefer_end=False), 'forecast': _month_to_date(forecast, prefer_end=True),
            'responsible': 'Roger Campbell /  delivery team'
        })
    red=''
    ma=re.search(r'Main deviations and mitigating actions(.*?)(?:1\. Achievement|Behind / Major concern|Status:)', block, re.I|re.S)
    if ma:
        red=normalize(ma.group(1)).replace('Discipline RED Activity Mitigating action','').strip()
    return {'year': 2026, 'week': week, 'period_label': f'2026 W{week:02d}' if week else '', 'file_name': fname, 'issue_date': issue_date, 'section_title': 'BOP | Distribution () – Activities and Progress', 'source_note': 'BOP Distribution  page; PB &  progress and 6-month key waypoints', 'engineering_mdr': eng, 'procurement_tg': proc, 'ugp_schedule': sched, 'waypoints': waypoints, 'red_activity_text': red}

# risk summary parse



def parse_engineering_scope_history(text, fname):
    """Parse scope-specific SWECO Engineering Management readiness history.

    The SWECO progress-by-discipline chart is embedded as an image in the PDFs,
    so text extraction does not expose the bar values consistently. For the
    historical reports bundled with this app, the values below were read from
    the chart pages and are stored so engineering progress is week-aware instead
    of locked to the latest report only.
    """
    week = None
    m = re.search(r'Weekly Report\s+W\s*(\d{1,2})', text, re.I)
    if m:
        week = int(m.group(1))
    known = {
        16: {'rail': (98.20, 'ETC 15 Apr 2026 / then ROS chart moved to 99.43% in W17'), 'ponds': (97.00, 'ETC late June 2026'), 'roads': (92.14, 'ETC July 2026')},
        17: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (97.00, 'ETC late June 2026'), 'roads': (90.96, 'ETC July 2026')},
        18: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (97.00, 'ETC late June 2026'), 'roads': (89.62, 'ETC July 2026')},
        19: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (98.50, 'ETC late June 2026'), 'roads': (89.62, 'ETC July 2026')},
        20: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (98.50, 'ETC late June 2026'), 'roads': (89.62, 'ETC July 2026')},
        21: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (98.50, 'ETC late June 2026'), 'roads': (89.62, 'ETC July 2026')},
        22: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (98.50, 'ETC late June 2026'), 'roads': (89.62, 'ETC July 2026')},
        23: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (98.50, 'ETC late June 2026'), 'roads': (89.62, 'ETC July 2026')},
        24: {'rail': (99.43, 'ETC 15 Apr 2026 / complete'), 'ponds': (98.50, 'ETC late June 2026'), 'roads': (89.62, 'ETC July 2026')},
        26: {'rail': (99.43, '100% ROS Sig'), 'ponds': (98.50, '1.50% remaining'), 'roads': (89.62, '10.38% remaining')},
    }
    page_by_week = {16: 26, 17: 26, 18: 28, 19: 27, 20: 24, 21: 24, 22: 23, 23: 23, 24: 23, 26: 22}
    names = {'rail': 'Rail', 'ponds': 'Ponds', 'roads': 'Roads'}
    rows = []
    if week in known:
        for scope_id, (val, etc) in known[week].items():
            rows.append({
                'year': 2026,
                'week': week,
                'scope_id': scope_id,
                'area': names[scope_id],
                'actual': val,
                'forecast': val,
                'deviation': 0.0,
                'etc': etc,
                'file_name': fname,
                'source_page': page_by_week.get(week),
                'basis': f'SWECO Engineering Management scope readiness chart, 2026 W{week:02d}',
            })
    return rows

def parse_risks(text, fname):
    week=None
    m=re.search(r'Weekly Report\s+W\s*(\d{1,2})', text, re.I)
    if m: week=int(m.group(1))
    summary={'week':week,'file_name':fname,'total_risks':None,'active_risks':None,'critical':None,'major':None,'moderate':None,'minor':None,'mit_total':None,'mit_overdue':None,'mit_open':None,'mit_closed':None,'mit_canceled':None}
    m=re.search(r'(\d+) total BOP risks identified\s*[–-]\s*(\d+) active', text, re.I)
    if m: summary['total_risks']=int(m.group(1)); summary['active_risks']=int(m.group(2))
    m=re.search(r'Out of \d+ active risk\s*[–-]\s*(\d+) critical,\s*(\d+) major,\s*(\d+) moderate,\s*(\d+) minor', text, re.I)
    if m:
        summary['critical']=int(m.group(1)); summary['major']=int(m.group(2)); summary['moderate']=int(m.group(3)); summary['minor']=int(m.group(4))
    m=re.search(r'(\d+) total mitigation actions planned.*?[–-]\s*(\d+) overdue,\s*(\d+) open,\s*(\d+) closed,\s*(\d+) canceled', text, re.I|re.S)
    if m:
        summary['mit_total']=int(m.group(1)); summary['mit_overdue']=int(m.group(2)); summary['mit_open']=int(m.group(3)); summary['mit_closed']=int(m.group(4)); summary['mit_canceled']=int(m.group(5))
    # top risks: find risk id blocks
    risk_rows=[]
    for rm in re.finditer(r'(B1\.BOP\.R\d+)\s+(.*?)\s+(\d{2})\s+(.*?)(\d{2})\s*(?=\n\s*B1\.BOP\.R|\Z)', text, re.I|re.S):
        rid=rm.group(1); mid=normalize(rm.group(2)); current=int(rm.group(3)); actions=normalize(rm.group(4)); pred=int(rm.group(5))
        # attempt title after function words contains hyphen
        title=mid
        if len(title)>220: title=title[:220]
        risk_rows.append({'week':week,'file_name':fname,'risk_id':rid,'risk_title':title,'current_score':current,'predicted_score':pred,'actions_excerpt':actions[:500]})
    return summary, risk_rows

texts={}
for f in files:
    txt=extract_text(f)
    texts[f.name]=txt
    out_path=DATA_DIR/f"extracted_{f.stem}.txt"
    out_path.write_text(txt, encoding='utf-8')

hse=[]; schedules=[]; waypoints=[]; risks_summary=[]; risks=[]
all_watch=[]
engineering_scope_history=[]
ugp_distribution_history=[]
for fname,txt in texts.items():
    hse.append(parse_hse(txt, fname))
    schedules.append(parse_overall_schedule(txt, fname))
    waypoints.extend(parse_waypoints(txt, fname))
    rs, rr=parse_risks(txt, fname); risks_summary.append(rs); risks.extend(rr)
# Build watchlist from latest report (max week) and some targeted older reports if desired
latest_week=max([r['week'] or 0 for r in hse])
for fname,txt in texts.items():
    # Keep all parsed watchlist history. The app uses year + week labels, so
    # historical uploads from 2024, 2025, 2026, etc. can be trended without
    # mixing same-numbered weeks across years.
    all_watch.extend(parse_area_statuses(txt, fname))
    engineering_scope_history.extend(parse_engineering_scope_history(txt, fname))
    ugp_rec = parse_ugp_distribution(txt, fname)
    if ugp_rec:
        ugp_distribution_history.append(ugp_rec)
        sched = ugp_rec.get('ugp_schedule') or {}
        actual = sched.get('actual'); fc = sched.get('forecast')
        if actual is not None and fc is not None:
            add_watch(all_watch, ugp_rec.get('week'), fname, '', f" distribution progress: actual {actual:.2f}%, forecast {fc:.2f}%, deviation {actual-fc:+.2f}% (plan {sched.get('plan'):.2f}%, next FC {sched.get('next_fc'):.2f}%).", category='Progress', severity='Info', owner='Roger Campbell /  delivery team', action='Track against BOP Distribution  schedule page.', status='Monitor')
        if ugp_rec.get('red_activity_text'):
            add_watch(all_watch, ugp_rec.get('week'), fname, '', ugp_rec.get('red_activity_text'), category='BOP Distribution red activity', severity='High', owner='Roger Campbell /  delivery team', action='Close the active BOP Distribution  red activity and confirm recovery in next weekly report.', status='Open')
        for wp in ugp_rec.get('waypoints', []):
            if wp.get('forecast'):
                add_watch(all_watch, ugp_rec.get('week'), fname, '', f"{wp.get('name')} - planned {wp.get('planned_raw') or 'TBC'} / forecast {wp.get('forecast_raw') or 'TBC'}", category='BOP Distribution milestone', severity='Medium', owner=wp.get('responsible'), action='Control via  milestone forecast timeline.', status='Monitor', due=wp.get('forecast'))
        eng = ugp_rec.get('engineering_mdr') or {}
        if eng.get('actual') is not None:
            actual_eng=eng.get('actual'); fc_eng=eng.get('forecast') if eng.get('forecast') is not None else actual_eng
            engineering_scope_history.append({'year':2026,'week':ugp_rec.get('week'),'scope_id':'','area':'','actual':actual_eng,'forecast':fc_eng,'deviation':round(actual_eng-fc_eng,2) if actual_eng is not None and fc_eng is not None else None,'etc': f"Next FC {eng.get('next_fc'):.2f}% from BOP Distribution PB &  Engineering (MDR)" if eng.get('next_fc') is not None else 'BOP Distribution PB &  Engineering (MDR)','file_name':fname,'source_page':'BOP Distribution ','basis': f"BOP Distribution  - PB &  Engineering (MDR), {ugp_rec.get('period_label')}"})
# Also include detailed area statuses if latest has it

# Sort data
hse=sorted(hse, key=lambda r:(r['week'] or 999))
schedules=sorted(schedules, key=lambda r:(r['week'] or 999))
risks_summary=sorted(risks_summary, key=lambda r:(r['week'] or 999))
risks=sorted(risks, key=lambda r: (-(r.get('current_score') or 0), r['risk_id']))
all_watch=sorted(all_watch, key=lambda r: (r['week'] or 0, r['area'], {'High':0,'Medium':1,'Info':2}.get(r['severity'],3)), reverse=True)

bundle={
    'hse':hse,
    'schedule':schedules,
    'waypoints':waypoints,
    'risks_summary':risks_summary,
    'risks':risks,
    'watchlist':all_watch,
    'engineering_scope_history':engineering_scope_history,
    'ugp_distribution_history':ugp_distribution_history,
    'source_files':[
        {
            'file_name':f.name,
            'suffix':f.suffix,
            'text_chars':len(texts[f.name]),
            'extracted_text':texts[f.name],
        }
        for f in files
    ],
}
Path(os.environ.get('BOP_BUNDLE_PATH', str(DATA_DIR/'weekly_data_bundle.json'))).write_text(json.dumps(bundle, ensure_ascii=False, indent=2), encoding='utf-8')
print(json.dumps({'hse_rows':len(hse),'schedules':len(schedules),'waypoints':len(waypoints),'risks_summary':len(risks_summary),'risks':len(risks),'watchlist':len(all_watch),'latest_week':latest_week,'ugp_distribution_history':len(ugp_distribution_history)}, indent=2))
# Print sample latest HSE and schedule
print(json.dumps(hse[-3:], ensure_ascii=False, indent=2)[:3000])
